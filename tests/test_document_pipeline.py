from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path
from xml.sax.saxutils import escape
from zipfile import ZIP_DEFLATED, ZipFile

from pypdf import PdfWriter
from pptx import Presentation

from podcast_anything_local.core.config import Settings
from podcast_anything_local.services.document_pipeline import MultimodalDocumentService


class _FakeDocumentProvider:
    def __init__(self) -> None:
        self.chunk_calls: list[tuple[int, int, int, str | None]] = []
        self.map_calls: int = 0
        self.plan_calls: int = 0
        self.plan_lengths: list[str] = []

    def summarize_pdf_chunk(
        self,
        *,
        pdf_bytes: bytes,
        filename: str,
        title: str | None,
        chunk_index: int,
        chunk_count: int,
        page_start: int,
        page_end: int,
        script_mode: str,
        supplemental_text: str | None = None,
    ) -> dict[str, object]:
        self.chunk_calls.append((chunk_index, page_start, page_end, supplemental_text))
        return {
            "page_start": page_start,
            "page_end": page_end,
            "summary": f"Summary for pages {page_start}-{page_end}",
            "key_points": [f"Point {chunk_index}"],
            "visual_elements": [f"Visual {chunk_index}"],
            "podcast_angles": [f"Angle {chunk_index}"],
            "must_include_details": [f"Detail {chunk_index}"],
            "caveats": [],
        }

    def build_document_map(
        self,
        *,
        chunk_summaries: list[dict[str, object]],
        title: str | None,
        script_mode: str,
    ) -> dict[str, object]:
        self.map_calls += 1
        return {
            "overall_summary": "Combined summary",
            "narrative_arc": ["Start", "Middle", "End"],
            "must_include": ["Key fact"],
            "supporting_details": ["Supporting detail"],
            "visual_takeaways": ["Important chart"],
            "caveats": [],
        }

    def build_podcast_plan(
        self,
        *,
        document_map: dict[str, object],
        title: str | None,
        script_mode: str,
        podcast_length: str = "medium",
    ) -> dict[str, object]:
        self.plan_calls += 1
        self.plan_lengths.append(podcast_length)
        return {
            "working_title": "Quantum Document Brief",
            "audience": "General technical audience",
            "angle": "Explain the big picture clearly",
            "intro": "Set up the topic fast.",
            "outro": "Close with one takeaway.",
            "must_include": ["Key fact"],
            "segments": [
                {
                    "name": "Foundations",
                    "purpose": "Explain the core idea",
                    "beats": ["Beat A", "Beat B"],
                    "source_pages": [1, 2],
                }
            ],
        }


def _build_settings(tmp_path: Path) -> Settings:
    data_dir = tmp_path / "data"
    return Settings(
        app_env="test",
        app_name="Podcast Anything Local Test",
        data_dir=data_dir,
        database_path=data_dir / "app.db",
        jobs_dir=data_dir / "jobs",
        web_extractor="auto",
        rewrite_style="podcast",
        openai_base_url="https://api.openai.com/v1",
        openai_api_key=None,
        openai_model="gpt-4o-mini",
        tts_provider="openai",
        elevenlabs_api_key=None,
        elevenlabs_model_id="eleven_multilingual_v2",
        elevenlabs_output_format="mp3_44100_128",
    )


def _build_docx_bytes(paragraphs: list[str]) -> bytes:
    document_xml = "".join(
        f"<w:p><w:r><w:t>{escape(paragraph)}</w:t></w:r></w:p>"
        for paragraph in paragraphs
    )
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>
"""
    relationships = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>
"""
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    {document_xml}
    <w:sectPr/>
  </w:body>
</w:document>
"""
    buffer = BytesIO()
    with ZipFile(buffer, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)
    return buffer.getvalue()


def test_multimodal_document_service_builds_chunks_and_plan(tmp_path: Path) -> None:
    pdf_path = tmp_path / "document.pdf"
    writer = PdfWriter()
    for _ in range(7):
        writer.add_blank_page(width=300, height=300)
    with pdf_path.open("wb") as file:
        writer.write(file)

    fake_provider = _FakeDocumentProvider()
    service = MultimodalDocumentService(
        _build_settings(tmp_path),
        provider_factory=lambda: fake_provider,
    )

    assert service.should_use(
        source_type="pdf",
        source_file_path=str(pdf_path),
    )

    analysis = service.analyze_pdf_document(
        source_file_path=str(pdf_path),
        source_display_name="document.pdf",
        title="Quantum Notes",
        script_mode="single",
    )
    assert analysis.page_count == 7
    assert fake_provider.chunk_calls == [(1, 1, 6, None), (2, 6, 7, None)]
    assert fake_provider.map_calls == 1

    podcast_plan = service.build_podcast_plan(
        document_map=analysis.document_map,
        title="Quantum Notes",
        script_mode="single",
        podcast_length="long",
    )
    rewrite_source = service.build_rewrite_source_text(
        analysis=analysis,
        podcast_plan=podcast_plan,
    )
    artifacts = service.build_artifacts(analysis=analysis)
    plan_artifacts = service.build_plan_artifacts(
        podcast_plan=podcast_plan,
        rewrite_source_text=rewrite_source,
    )

    assert fake_provider.plan_calls == 1
    assert fake_provider.plan_lengths == ["long"]
    assert "Quantum Document Brief" in rewrite_source
    assert "Chunk evidence:" in rewrite_source
    assert "page_index.json" in artifacts
    assert "document_map.json" in artifacts
    assert "chunk_001_summary.json" in artifacts
    assert "podcast_plan.json" in plan_artifacts
    assert "rewrite_input.txt" in plan_artifacts


def test_prepare_document_for_analysis_normalizes_docx_to_pdf(tmp_path: Path) -> None:
    docx_path = tmp_path / "brief.docx"
    docx_path.write_bytes(
        _build_docx_bytes(
            [
                "Quantum mechanics explains microscopic systems.",
                "Its applications include semiconductors and lasers.",
            ]
        )
    )

    service = MultimodalDocumentService(_build_settings(tmp_path))
    source_text = (
        "Quantum mechanics explains microscopic systems.\n\n"
        "Its applications include semiconductors and lasers."
    )
    prepared = service.prepare_document_for_analysis(
        source_text=source_text,
        source_type="docx",
        source_file_path=str(docx_path),
        source_file_name="brief.docx",
    )

    assert prepared.analysis_pdf_bytes is not None
    assert prepared.analysis_artifact_name == "normalized.pdf"
    assert prepared.metadata["normalized_document_used"] is True
    assert prepared.metadata["normalized_document_source_type"] == "docx"
    assert prepared.page_context
    assert "normalized_page_context.json" in prepared.text_artifacts


def test_prepare_document_for_analysis_normalizes_pptx_and_feeds_notes_into_chunk_prompt(
    tmp_path: Path
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quantum basics"
    slide.placeholders[1].text = "Wave functions\nUncertainty principle"
    slide.notes_slide.notes_text_frame.text = "Introduce the topic with simple examples."
    buffer = BytesIO()
    presentation.save(buffer)
    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(buffer.getvalue())

    source_text = (
        "Slide 1: Quantum basics\n\n"
        "Wave functions\n\n"
        "Uncertainty principle\n\n"
        "Speaker notes:\n"
        "Introduce the topic with simple examples."
    )
    fake_provider = _FakeDocumentProvider()
    service = MultimodalDocumentService(
        _build_settings(tmp_path),
        provider_factory=lambda: fake_provider,
    )
    prepared = service.prepare_document_for_analysis(
        source_text=source_text,
        source_type="pptx",
        source_file_path=str(pptx_path),
        source_file_name="deck.pptx",
    )

    normalized_pdf_path = tmp_path / "normalized.pdf"
    normalized_pdf_path.write_bytes(prepared.analysis_pdf_bytes or b"")

    analysis = service.analyze_pdf_document(
        source_file_path=str(normalized_pdf_path),
        source_display_name="deck.pdf",
        title="Quantum basics",
        script_mode="single",
        page_context=prepared.page_context,
    )

    slide_notes = json.loads(prepared.text_artifacts["slide_notes.json"])
    assert prepared.analysis_pdf_bytes is not None
    assert prepared.metadata["normalized_document_has_slide_notes"] is True
    assert slide_notes[0]["notes"] == "Introduce the topic with simple examples."
    assert analysis.page_count >= 1
    assert "Introduce the topic with simple examples." in (fake_provider.chunk_calls[0][3] or "")
