from __future__ import annotations

import json
import time
from io import BytesIO
from pathlib import Path

from fastapi.testclient import TestClient
from pypdf import PdfWriter
from pptx import Presentation

from podcast_anything_local.core.config import Settings
from podcast_anything_local.main import create_app
from podcast_anything_local.services.ingestion import IngestionService


def _build_settings(tmp_path: Path) -> Settings:
    data_dir = tmp_path / "data"
    return Settings(
        app_env="test",
        app_name="Podcast Anything Local Test",
        data_dir=data_dir,
        database_path=data_dir / "app.db",
        jobs_dir=data_dir / "jobs",
        web_extractor="auto",
        rewrite_provider="demo",
        rewrite_style="podcast",
        ollama_base_url="http://localhost:11434/api",
        ollama_model="gemma3:4b",
        ollama_generate_timeout_seconds=600,
        openai_base_url="https://api.openai.com/v1",
        openai_api_key=None,
        openai_model="gpt-4o-mini",
        tts_provider="wave",
        tts_default_voice="host_a",
        tts_duo_voice="host_b",
        piper_model_path=None,
        piper_model_path_b=None,
        piper_config_path=None,
        piper_config_path_b=None,
        piper_speaker_id=None,
        piper_speaker_id_b=None,
        elevenlabs_api_key=None,
        elevenlabs_model_id="eleven_multilingual_v2",
        elevenlabs_output_format="mp3_44100_128",
    )


def _build_openai_settings(tmp_path: Path) -> Settings:
    data_dir = tmp_path / "data"
    return Settings(
        app_env="test",
        app_name="Podcast Anything Local Test",
        data_dir=data_dir,
        database_path=data_dir / "app.db",
        jobs_dir=data_dir / "jobs",
        web_extractor="auto",
        rewrite_provider="openai",
        rewrite_style="podcast",
        ollama_base_url="http://localhost:11434/api",
        ollama_model="gemma3:4b",
        ollama_generate_timeout_seconds=600,
        openai_base_url="https://api.openai.com/v1",
        openai_api_key="test-key",
        openai_model="gpt-4o-mini",
        tts_provider="wave",
        tts_default_voice="host_a",
        tts_duo_voice="host_b",
        piper_model_path=None,
        piper_model_path_b=None,
        piper_config_path=None,
        piper_config_path_b=None,
        piper_speaker_id=None,
        piper_speaker_id_b=None,
        elevenlabs_api_key=None,
        elevenlabs_model_id="eleven_multilingual_v2",
        elevenlabs_output_format="mp3_44100_128",
    )


def _wait_for_terminal_job_state(client: TestClient, job_id: str) -> dict:
    deadline = time.time() + 5
    while time.time() < deadline:
        response = client.get(f"/jobs/{job_id}")
        payload = response.json()
        if payload["status"] in {"completed", "failed"}:
            return payload
        time.sleep(0.1)
    raise AssertionError("Job did not reach a terminal state within 5 seconds.")


def test_create_job_from_url_runs_pipeline(tmp_path: Path, monkeypatch) -> None:
    source_text = "Example source material for a local podcast run."

    def fake_ingest(self, *, source_kind=None, source_url=None, source_file_path=None, source_file_name=None):
        assert source_kind == "url"
        assert source_url == "https://example.com/article"
        return (
            source_text,
            {"source_type": "webpage", "source_char_count": len(source_text)},
        )

    monkeypatch.setattr(IngestionService, "ingest", fake_ingest)

    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            json={
                "source_url": "https://example.com/article",
                "script_mode": "duo",
            },
        )

        assert response.status_code == 202
        payload = response.json()
        assert payload["status"] == "queued"
        assert payload["script_mode"] == "duo"

        terminal_payload = _wait_for_terminal_job_state(client, payload["job_id"])
        assert terminal_payload["status"] == "completed"
        assert terminal_payload["title"] == "Local Podcast Conversation"
        assert terminal_payload["audio_artifact"].endswith("audio.wav")
        assert terminal_payload["metadata"]["rewrite_input_char_count"] == len(source_text)
        assert terminal_payload["metadata"]["rewrite_input_truncated"] is False
        assert terminal_payload["metadata"]["title_source"] == "llm"

        artifacts_response = client.get(f"/jobs/{payload['job_id']}/artifacts")
        artifacts_payload = artifacts_response.json()
        artifacts = {item["name"] for item in artifacts_payload}
        assert {"audio.wav", "metadata.json", "script.txt", "source.txt"} <= artifacts
        script_artifact = next(item for item in artifacts_payload if item["name"] == "script.txt")
        assert script_artifact["download_path"] == f"/jobs/{payload['job_id']}/artifacts/script.txt"

        download_response = client.get(script_artifact["download_path"])
        assert download_response.status_code == 200
        assert "Welcome back" in download_response.text or len(download_response.text) > 20


def test_create_job_from_uploaded_txt_file(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={"script_mode": "single"},
            files={"source_file": ("brief.txt", b"Uploaded notes for a short podcast.")},
        )

        assert response.status_code == 202
        payload = response.json()
        assert payload["source_kind"] == "file"
        assert payload["source_file_name"] == "brief.txt"

        terminal_payload = _wait_for_terminal_job_state(client, payload["job_id"])
        assert terminal_payload["status"] == "completed"
        assert terminal_payload["title"] == "Local Podcast Draft"
        assert terminal_payload["source_artifact"].endswith("source.txt")
        assert terminal_payload["metadata"]["rewrite_input_truncated"] is False
        assert terminal_payload["metadata"]["title_source"] == "llm"

        artifacts_response = client.get(f"/jobs/{payload['job_id']}/artifacts")
        artifacts_payload = artifacts_response.json()
        artifacts = {item["name"] for item in artifacts_payload}
        assert {"audio.wav", "input_brief.txt", "metadata.json", "script.txt", "source.txt"} <= artifacts

        audio_response = client.get(f"/jobs/{payload['job_id']}/artifacts/audio.wav")
        assert audio_response.status_code == 200
        assert audio_response.headers["content-type"] in {"audio/wav", "audio/x-wav"}
        assert len(audio_response.content) > 100


def test_create_job_from_pasted_text(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={
                "source_text": "Pasted text for a short podcast draft.",
                "script_mode": "single",
            },
        )

        assert response.status_code == 202
        payload = response.json()
        assert payload["source_kind"] == "text"
        assert payload["source_file_name"] == "pasted_text.txt"

        terminal_payload = _wait_for_terminal_job_state(client, payload["job_id"])
        assert terminal_payload["status"] == "completed"
        assert terminal_payload["metadata"]["source_type"] == "text"
        assert terminal_payload["metadata"]["submitted_text_char_count"] == 38

        artifacts_response = client.get(f"/jobs/{payload['job_id']}/artifacts")
        artifacts_payload = artifacts_response.json()
        artifacts = {item["name"] for item in artifacts_payload}
        assert {"audio.wav", "input_pasted_text.txt", "metadata.json", "script.txt", "source.txt"} <= artifacts


def test_create_job_from_uploaded_pptx_file(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quantum deck"
    slide.placeholders[1].text = "Quanta\nPhotons"
    slide.notes_slide.notes_text_frame.text = "Frame this as a short primer."
    buffer = BytesIO()
    presentation.save(buffer)

    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={"script_mode": "single"},
            files={
                "source_file": (
                    "deck.pptx",
                    buffer.getvalue(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

        assert response.status_code == 202
        payload = response.json()
        assert payload["source_kind"] == "file"
        assert payload["source_file_name"] == "deck.pptx"

        terminal_payload = _wait_for_terminal_job_state(client, payload["job_id"])
        assert terminal_payload["status"] == "completed"
        assert terminal_payload["metadata"]["source_type"] == "pptx"

        artifacts_response = client.get(f"/jobs/{payload['job_id']}/artifacts")
        artifacts_payload = artifacts_response.json()
        artifacts = {item["name"] for item in artifacts_payload}
        assert {"audio.wav", "input_deck.pptx", "metadata.json", "script.txt", "source.txt"} <= artifacts


def test_create_job_from_uploaded_pdf_uses_multimodal_document_pipeline(
    tmp_path: Path, monkeypatch
) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    writer.add_blank_page(width=300, height=300)
    buffer = BytesIO()
    writer.write(buffer)

    def fake_post(url, headers=None, json=None, timeout=0):
        if url.endswith("/responses"):
            prompt_text = json["input"][0]["content"][0]["text"]
            if "chunk-level document summaries" in prompt_text:
                payload = {
                    "output": [
                        {
                            "content": [
                                {
                                    "type": "output_text",
                                    "text": json_module_dumps(
                                        {
                                            "overall_summary": "Mapped summary",
                                            "narrative_arc": ["Arc start", "Arc end"],
                                            "must_include": ["Must include"],
                                            "supporting_details": ["Supporting detail"],
                                            "visual_takeaways": ["Visual takeaway"],
                                            "caveats": [],
                                        }
                                    ),
                                }
                            ]
                        }
                    ]
                }
            elif "structured document map" in prompt_text:
                payload = {
                    "output": [
                        {
                            "content": [
                                {
                                    "type": "output_text",
                                    "text": json_module_dumps(
                                        {
                                            "working_title": "Quantum PDF Brief",
                                            "audience": "General audience",
                                            "angle": "Explain the essentials",
                                            "intro": "Start quickly.",
                                            "outro": "End clearly.",
                                            "must_include": ["Must include"],
                                            "segments": [
                                                {
                                                    "name": "Segment One",
                                                    "purpose": "Explain the core concept",
                                                    "beats": ["Beat one", "Beat two"],
                                                    "source_pages": [1, 2],
                                                }
                                            ],
                                        }
                                    ),
                                }
                            ]
                        }
                    ]
                }
            else:
                payload = {
                    "output": [
                        {
                            "content": [
                                {
                                    "type": "output_text",
                                    "text": json_module_dumps(
                                        {
                                            "page_start": 1,
                                            "page_end": 2,
                                            "summary": "Chunk summary",
                                            "key_points": ["Point"],
                                            "visual_elements": ["Visual"],
                                            "podcast_angles": ["Angle"],
                                            "must_include_details": ["Detail"],
                                            "caveats": [],
                                        }
                                    ),
                                }
                            ]
                        }
                    ]
                }
            return _FakeResponse(payload)

        if url.endswith("/chat/completions"):
            prompt = json["messages"][0]["content"]
            if "You are titling a podcast episode." in prompt:
                return _FakeResponse({"choices": [{"message": {"content": "Title: Quantum PDF Brief"}}]})
            return _FakeResponse(
                {
                    "choices": [
                        {
                            "message": {
                                "content": "Welcome back. Today we're turning a PDF into a podcast."
                            }
                        }
                    ]
                }
            )

        raise AssertionError(f"Unexpected URL: {url}")

    class _FakeResponse:
        def __init__(self, payload: dict, status_code: int = 200) -> None:
            self._payload = payload
            self.status_code = status_code
            self.headers: dict[str, str] = {}

        def raise_for_status(self) -> None:
            if self.status_code >= 400:
                raise RuntimeError(self.status_code)

        def json(self) -> dict:
            return self._payload

    def json_module_dumps(payload: dict) -> str:
        return json.dumps(payload, ensure_ascii=True)

    monkeypatch.setattr("requests.post", fake_post)

    app = create_app(_build_openai_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={"script_mode": "single"},
            files={"source_file": ("brief.pdf", buffer.getvalue(), "application/pdf")},
        )

        assert response.status_code == 202
        payload = response.json()

        terminal_payload = _wait_for_terminal_job_state(client, payload["job_id"])
        assert terminal_payload["status"] == "completed"
        assert terminal_payload["metadata"]["multimodal_document_pipeline"] is True
        assert terminal_payload["metadata"]["multimodal_document_page_count"] == 2
        assert terminal_payload["metadata"]["multimodal_document_chunk_count"] == 1
        assert terminal_payload["title"] == "Quantum PDF Brief"

        artifacts = {
            item["name"] for item in client.get(f"/jobs/{payload['job_id']}/artifacts").json()
        }
        assert {
            "audio.wav",
            "chunk_001_summary.json",
            "document_map.json",
            "input_brief.pdf",
            "metadata.json",
            "page_index.json",
            "podcast_plan.json",
            "rewrite_input.txt",
            "script.txt",
            "source.txt",
        } <= artifacts


def test_download_job_artifact_returns_404_for_missing_artifact(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={"script_mode": "single"},
            files={"source_file": ("brief.txt", b"Uploaded notes for a short podcast.")},
        )

        job_id = response.json()["job_id"]
        _wait_for_terminal_job_state(client, job_id)

        missing_response = client.get(f"/jobs/{job_id}/artifacts/does-not-exist.txt")
        assert missing_response.status_code == 404


def test_retry_job_clears_stale_generated_artifacts(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        response = client.post(
            "/jobs",
            data={"script_mode": "single"},
            files={"source_file": ("brief.txt", b"Uploaded notes for a short podcast.")},
        )

        job_id = response.json()["job_id"]
        _wait_for_terminal_job_state(client, job_id)

        artifacts_before_retry = {
            item["name"] for item in client.get(f"/jobs/{job_id}/artifacts").json()
        }
        assert {"audio.wav", "input_brief.txt", "metadata.json", "script.txt", "source.txt"} <= artifacts_before_retry

        client.app.state.executor.submit = lambda _: None
        retry_response = client.post(f"/jobs/{job_id}/retry")
        assert retry_response.status_code == 200

        retry_payload = retry_response.json()
        assert retry_payload["status"] == "queued"
        assert retry_payload["current_stage"] is None
        assert retry_payload["title"] is None
        assert retry_payload["source_artifact"] is None
        assert retry_payload["script_artifact"] is None
        assert retry_payload["audio_artifact"] is None
        assert retry_payload["error"] is None
        assert retry_payload["metadata"]["retry_count"] == 1
        assert "audio_file_name" not in retry_payload["metadata"]
        assert "metadata_artifact" not in retry_payload["metadata"]
        assert "script_char_count" not in retry_payload["metadata"]
        assert "source_char_count" not in retry_payload["metadata"]
        assert "source_type" not in retry_payload["metadata"]
        assert "title_source" not in retry_payload["metadata"]

        artifacts_after_retry = {
            item["name"] for item in client.get(f"/jobs/{job_id}/artifacts").json()
        }
        assert artifacts_after_retry == {"input_brief.txt"}


def test_rewrite_stream_route_replays_snapshot_and_completion(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        client.app.state.executor.submit = lambda _: None
        response = client.post(
            "/jobs",
            json={
                "source_url": "https://example.com/article",
                "script_mode": "single",
            },
        )

        job_id = response.json()["job_id"]
        broker = client.app.state.job_event_broker
        broker.publish_rewrite_chunk(job_id, "Hello ")
        broker.publish_rewrite_chunk(job_id, "world.")
        broker.publish_rewrite_complete(job_id)

        with client.stream("GET", f"/jobs/{job_id}/rewrite-stream") as stream_response:
            body = "".join(stream_response.iter_text())

        assert stream_response.status_code == 200
        assert "event: rewrite_snapshot" in body
        assert '"text": "Hello world."' in body
        assert "event: rewrite_complete" in body


def test_rewrite_preview_route_returns_current_snapshot(tmp_path: Path) -> None:
    app = create_app(_build_settings(tmp_path))
    with TestClient(app) as client:
        client.app.state.executor.submit = lambda _: None
        response = client.post(
            "/jobs",
            json={
                "source_url": "https://example.com/article",
                "script_mode": "single",
            },
        )

        job_id = response.json()["job_id"]
        client.app.state.job_event_broker.publish_rewrite_chunk(job_id, "Partial script")

        preview_response = client.get(f"/jobs/{job_id}/rewrite-preview")

        assert preview_response.status_code == 200
        assert preview_response.json() == {"text": "Partial script"}
