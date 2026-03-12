from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from podcast_anything_local.services.ingestion import (
    IngestionService,
    IngestionServiceError,
    _extract_article_text,
    _extract_document_text,
)


def test_extract_article_text_uses_bs4_when_selected(monkeypatch) -> None:
    monkeypatch.setattr(
        "podcast_anything_local.services.ingestion.trafilatura.extract",
        lambda *args, **kwargs: "Trafilatura text",
    )

    html = """
    <html>
      <body>
        <article>
          <p>BeautifulSoup paragraph one.</p>
          <p>BeautifulSoup paragraph two.</p>
        </article>
      </body>
    </html>
    """

    extracted = _extract_article_text(html, web_extractor="bs4")

    assert extracted == "BeautifulSoup paragraph one.\nBeautifulSoup paragraph two."


def test_extract_article_text_falls_back_to_bs4_in_auto_mode(monkeypatch) -> None:
    monkeypatch.setattr(
        "podcast_anything_local.services.ingestion.trafilatura.extract",
        lambda *args, **kwargs: None,
    )

    html = """
    <html>
      <body>
        <article>
          <p>Fallback paragraph.</p>
        </article>
      </body>
    </html>
    """

    extracted = _extract_article_text(html, web_extractor="auto")

    assert extracted == "Fallback paragraph."


def test_extract_article_text_errors_when_trafilatura_selected_and_empty(monkeypatch) -> None:
    monkeypatch.setattr(
        "podcast_anything_local.services.ingestion.trafilatura.extract",
        lambda *args, **kwargs: None,
    )

    try:
        _extract_article_text("<html><body><p>Ignored</p></body></html>", web_extractor="trafilatura")
    except IngestionServiceError as exc:
        assert "trafilatura" in str(exc)
    else:
        raise AssertionError("Expected IngestionServiceError")


def test_ingestion_service_validates_web_extractor() -> None:
    try:
        IngestionService(web_extractor="invalid")
    except IngestionServiceError as exc:
        assert "web_extractor" in str(exc)
    else:
        raise AssertionError("Expected IngestionServiceError")


def test_extract_document_text_supports_pptx_with_slide_notes() -> None:
    presentation = Presentation()

    slide_one = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide_one.shapes.title.text = "Quantum basics"
    slide_one.placeholders[1].text = "Wave functions\nUncertainty principle"
    slide_one.notes_slide.notes_text_frame.text = "Introduce the topic with simple examples."

    slide_two = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide_two.shapes.title.text = "Applications"
    slide_two.placeholders[1].text = "Semiconductors\nLasers"

    buffer = BytesIO()
    presentation.save(buffer)

    extracted_text, source_type = _extract_document_text(buffer.getvalue(), "deck.pptx")

    assert source_type == "pptx"
    assert "Slide 1: Quantum basics" in extracted_text
    assert "Wave functions" in extracted_text
    assert "Speaker notes:" in extracted_text
    assert "Introduce the topic with simple examples." in extracted_text
    assert "Slide 2: Applications" in extracted_text
    assert "Semiconductors" in extracted_text
