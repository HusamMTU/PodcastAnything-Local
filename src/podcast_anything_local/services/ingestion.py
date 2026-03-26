"""Input ingestion for webpages, YouTube URLs, and uploaded documents."""

from __future__ import annotations

import re
from collections.abc import Iterable
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from types import SimpleNamespace
from urllib.parse import parse_qs, urlparse

import requests
import trafilatura
from bs4 import BeautifulSoup

try:
    import docx2txt
except ModuleNotFoundError:  # pragma: no cover
    docx2txt = SimpleNamespace(process=None)

try:
    from pypdf import PdfReader
except ModuleNotFoundError:  # pragma: no cover
    PdfReader = None

try:
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover
    Presentation = None


class IngestionServiceError(RuntimeError):
    """Raised when source ingestion fails."""


_SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
}

_YOUTUBE_HOSTS = {
    "youtube.com",
    "www.youtube.com",
    "m.youtube.com",
    "youtu.be",
    "www.youtu.be",
}


class IngestionService:
    def __init__(self, *, web_extractor: str = "auto") -> None:
        normalized = web_extractor.strip().lower()
        if normalized not in {"auto", "trafilatura", "bs4"}:
            raise IngestionServiceError("web_extractor must be one of: auto, trafilatura, bs4")
        self._web_extractor = normalized

    def ingest(
        self,
        *,
        source_kind: str | None = None,
        source_url: str | None = None,
        source_file_path: str | None = None,
        source_file_name: str | None = None,
    ) -> tuple[str, dict[str, object]]:
        if bool(source_url) == bool(source_file_path):
            raise IngestionServiceError(
                "Provide exactly one of source_url or source_file_path for ingestion."
            )

        if source_url:
            if _is_youtube_url(source_url):
                text = _fetch_youtube_transcript(source_url)
                return text, {"source_type": "youtube", "source_char_count": len(text)}
            html = _fetch_html(source_url)
            text = _extract_article_text(html, web_extractor=self._web_extractor)
            return text, {"source_type": "webpage", "source_char_count": len(text)}

        assert source_file_path is not None
        file_path = Path(source_file_path)
        if not file_path.is_file():
            raise IngestionServiceError(f"Uploaded file not found: {source_file_path}")
        if source_kind == "text":
            text = _extract_txt_text(file_path.read_bytes())
            if not text:
                raise IngestionServiceError("No readable text found in pasted text input.")
            return text, {"source_type": "text", "source_char_count": len(text)}
        text, source_type = _extract_document_text(
            file_bytes=file_path.read_bytes(),
            filename=source_file_name or file_path.name,
        )
        return text, {"source_type": source_type, "source_char_count": len(text)}


def _fetch_html(url: str, timeout_sec: int = 30) -> str:
    if not url.startswith(("http://", "https://")):
        raise IngestionServiceError("source_url must start with http:// or https://")
    try:
        response = requests.get(
            url,
            timeout=timeout_sec,
            headers={"User-Agent": "podcast-anything-local/0.1"},
        )
        response.raise_for_status()
    except requests.RequestException as exc:
        raise IngestionServiceError(f"Failed to fetch webpage: {exc}") from exc
    return response.text


def _clean_text(lines: Iterable[str]) -> str:
    joined = "\n".join(line.strip() for line in lines if line and line.strip())
    return re.sub(r"\n{3,}", "\n\n", joined).strip()


def _extract_article_text(html: str, *, web_extractor: str = "auto") -> str:
    if web_extractor == "trafilatura":
        return _extract_article_text_with_trafilatura(html)
    if web_extractor == "bs4":
        return _extract_article_text_with_bs4(html)

    try:
        return _extract_article_text_with_trafilatura(html)
    except IngestionServiceError:
        return _extract_article_text_with_bs4(html)


def _extract_article_text_with_trafilatura(html: str) -> str:
    extracted = trafilatura.extract(
        html,
        include_comments=False,
        include_tables=False,
        output_format="txt",
    )
    if extracted:
        cleaned = _clean_text(extracted.splitlines())
        if cleaned:
            return cleaned
    raise IngestionServiceError("No readable text found in webpage with trafilatura.")


def _extract_article_text_with_bs4(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    article = soup.find("article")
    if article:
        paragraphs = [p.get_text(" ", strip=True) for p in article.find_all("p")]
        cleaned = _clean_text(paragraphs)
        if cleaned:
            return cleaned

    paragraphs = [p.get_text(" ", strip=True) for p in soup.find_all("p")]
    cleaned = _clean_text(paragraphs)
    if not cleaned:
        raise IngestionServiceError("No readable text found in webpage.")
    return cleaned


def _detect_document_type(filename: str) -> str:
    extension = Path(filename).suffix.lower()
    document_type = _SUPPORTED_EXTENSIONS.get(extension)
    if not document_type:
        supported = ", ".join(sorted(_SUPPORTED_EXTENSIONS))
        raise IngestionServiceError(
            f"Unsupported document type '{extension or filename}'. Use: {supported}"
        )
    return document_type


def _extract_document_text(file_bytes: bytes, filename: str) -> tuple[str, str]:
    if not file_bytes:
        raise IngestionServiceError("Uploaded document is empty.")
    document_type = _detect_document_type(filename)
    if document_type == "pdf":
        text = _extract_pdf_text(file_bytes)
    elif document_type == "docx":
        text = _extract_docx_text(file_bytes)
    elif document_type == "pptx":
        text = _extract_pptx_text(file_bytes)
    else:
        text = _extract_txt_text(file_bytes)
    if document_type == "pdf":
        return text, document_type
    if not text:
        raise IngestionServiceError(f"No readable text found in uploaded {document_type}.")
    return text, document_type


def _extract_pdf_text(file_bytes: bytes) -> str:
    if PdfReader is None:
        raise IngestionServiceError("PDF support requires the `pypdf` package.")
    try:
        reader = PdfReader(BytesIO(file_bytes))
    except Exception as exc:  # pragma: no cover
        raise IngestionServiceError(f"Failed to parse PDF document: {exc}") from exc
    extracted: list[str] = []
    for page in reader.pages:
        page_text = (page.extract_text() or "").strip()
        if page_text:
            extracted.append(page_text)
    return "\n\n".join(extracted).strip()


def _extract_docx_text(file_bytes: bytes) -> str:
    if docx2txt.process is None:
        raise IngestionServiceError("DOCX support requires the `docx2txt` package.")
    with NamedTemporaryFile(suffix=".docx") as temp_file:
        temp_file.write(file_bytes)
        temp_file.flush()
        try:
            raw_text = docx2txt.process(temp_file.name) or ""
        except Exception as exc:  # pragma: no cover
            raise IngestionServiceError(f"Failed to parse DOCX document: {exc}") from exc
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    return "\n\n".join(lines).strip()


def _extract_pptx_text(file_bytes: bytes) -> str:
    if Presentation is None:
        raise IngestionServiceError("PPTX support requires the `python-pptx` package.")
    try:
        presentation = Presentation(BytesIO(file_bytes))
    except Exception as exc:  # pragma: no cover
        raise IngestionServiceError(f"Failed to parse PPTX document: {exc}") from exc

    extracted: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_text = _shape_text(title_shape)
        slide_header = f"Slide {slide_number}"
        if title_text:
            slide_header = f"{slide_header}: {title_text}"

        body_blocks: list[str] = []
        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            text = _shape_text(shape)
            if text:
                body_blocks.append(text)

        notes_text = ""
        try:
            notes_text = _clean_text(slide.notes_slide.notes_text_frame.text.splitlines())
        except Exception:  # pragma: no cover
            notes_text = ""

        slide_parts = [slide_header]
        if body_blocks:
            slide_parts.append("\n\n".join(body_blocks))
        if notes_text:
            slide_parts.append(f"Speaker notes:\n{notes_text}")

        if len(slide_parts) > 1:
            extracted.append("\n\n".join(slide_parts))

    return "\n\n".join(extracted).strip()


def _shape_text(shape: object) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    raw_text = getattr(shape, "text", "") or ""
    return _clean_text(raw_text.splitlines())


def _extract_txt_text(file_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "latin-1"):
        try:
            return file_bytes.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise IngestionServiceError("Failed to decode TXT document with supported encodings.")


def _is_youtube_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
    except Exception:
        return False
    return (parsed.hostname or "").lower() in _YOUTUBE_HOSTS


def _extract_video_id(url: str) -> str:
    parsed = urlparse(url)
    host = (parsed.hostname or "").lower()
    if host not in _YOUTUBE_HOSTS:
        raise IngestionServiceError("source_url is not a supported YouTube URL")

    if host in {"youtu.be", "www.youtu.be"}:
        video_id = parsed.path.lstrip("/").split("/")[0]
        if video_id:
            return video_id

    if host in {"youtube.com", "www.youtube.com", "m.youtube.com"}:
        if parsed.path == "/watch":
            video_id = parse_qs(parsed.query).get("v", [None])[0]
            if video_id:
                return video_id
        for prefix in ("/shorts/", "/embed/", "/live/"):
            if parsed.path.startswith(prefix):
                video_id = parsed.path[len(prefix) :].split("/")[0]
                if video_id:
                    return video_id

    raise IngestionServiceError("Could not extract video id from YouTube URL")


def _fetch_youtube_transcript(url: str) -> str:
    video_id = _extract_video_id(url)
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except Exception as exc:  # pragma: no cover
        raise IngestionServiceError("youtube-transcript-api is not installed.") from exc

    try:
        if hasattr(YouTubeTranscriptApi, "get_transcript"):
            segments = YouTubeTranscriptApi.get_transcript(video_id, languages=["en", "en-US"])
        else:
            api = YouTubeTranscriptApi()
            try:
                segments = api.fetch(video_id, languages=["en", "en-US"])
            except TypeError:
                segments = api.fetch(video_id)
    except Exception as exc:
        raise IngestionServiceError(f"Failed to fetch YouTube transcript: {exc}") from exc

    if hasattr(segments, "to_raw_data"):
        segments = segments.to_raw_data()
    lines: list[str] = []
    for item in segments:
        text = item.get("text") if isinstance(item, dict) else getattr(item, "text", None)
        if isinstance(text, str) and text.strip():
            lines.append(re.sub(r"\s+", " ", text).strip())

    if not lines:
        raise IngestionServiceError("YouTube transcript is empty.")
    paragraphs: list[str] = []
    chunk: list[str] = []
    for line in lines:
        chunk.append(line)
        if len(chunk) >= 6:
            paragraphs.append(" ".join(chunk))
            chunk = []
    if chunk:
        paragraphs.append(" ".join(chunk))
    return "\n\n".join(paragraphs).strip()
